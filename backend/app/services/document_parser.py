from __future__ import annotations

from pathlib import Path

import fitz
from docx import Document
from pptx import Presentation


SUPPORTED_EXTENSIONS = {".pdf", ".docx", ".pptx", ".txt", ".md"}


def extract_text(path: str) -> str:
    """从支持的课程资料文件中提取纯文本。"""

    file_path = Path(path)
    suffix = file_path.suffix.lower()

    if suffix == ".pdf":
        return _extract_pdf(file_path)
    if suffix == ".docx":
        return _extract_docx(file_path)
    if suffix == ".pptx":
        return _extract_pptx(file_path)
    if suffix in {".txt", ".md"}:
        return file_path.read_text(encoding="utf-8", errors="ignore")

    raise ValueError(
        f"Unsupported file type: {suffix}. Supported: PDF, DOCX, PPTX, TXT, MD."
    )


def _extract_pdf(path: Path) -> str:
    """使用 PyMuPDF 提取每一页 PDF 文本。"""

    parts = []
    with fitz.open(path) as doc:
        for page_index, page in enumerate(doc, start=1):
            text = page.get_text("text").strip()
            if text:
                parts.append(f"[PDF page {page_index}]\n{text}")
    return "\n\n".join(parts)


def _extract_docx(path: Path) -> str:
    """从 DOCX 文件中提取段落和表格文本。"""

    document = Document(path)
    parts = [paragraph.text.strip() for paragraph in document.paragraphs]

    # Tables often contain exam points or formulas; include them as rows.
    for table in document.tables:
        for row in table.rows:
            cells = [cell.text.strip() for cell in row.cells if cell.text.strip()]
            if cells:
                parts.append(" | ".join(cells))

    return "\n\n".join(item for item in parts if item)


def _extract_pptx(path: Path) -> str:
    """从 PPTX 每一页幻灯片中提取可见文本。"""

    presentation = Presentation(path)
    parts = []
    for slide_index, slide in enumerate(presentation.slides, start=1):
        slide_parts = []
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text:
                slide_parts.append(shape.text.strip())
        if slide_parts:
            parts.append(f"[PPT slide {slide_index}]\n" + "\n".join(slide_parts))
    return "\n\n".join(parts)
